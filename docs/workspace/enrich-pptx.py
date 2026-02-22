"""
Enrich meeting-allineamento-vs.pptx with site screenshots.
Adds 4 new slides after relevant existing slides showing actual site pages
with annotations about team member contributions.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image
import copy, os

BASE = "/Users/luckyluke/projects/active/visione_sostenibile"
SCREENSHOTS = f"{BASE}/site-screenshots"
SRC = f"{BASE}/meeting-allineamento-vs.pptx"
OUT = f"{BASE}/workspace/meeting-enriched.pptx"

# Slide dimensions (16:9): 9144000 x 5143500 EMU = 10" x 5.63"
SLIDE_W = 9144000
SLIDE_H = 5143500

# Colors
DEEP_FOREST = RGBColor(0x0B, 0x1E, 0x0E)
PAPER_CANVAS = RGBColor(0xF2, 0xF0, 0xEC)
SUN_ACCENT = RGBColor(0xEA, 0xB8, 0x31)
LEAF_GREEN = RGBColor(0x4F, 0xA4, 0x5A)
LIGHT_GREEN = RGBColor(0x9C, 0xB8, 0x9F)
DARK_GREEN_BG = RGBColor(0x0F, 0x2B, 0x13)

def add_bg(slide, color):
    """Set solid background color on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, font_size=12, bold=False, italic=False, color=PAPER_CANVAS, align=PP_ALIGN.LEFT, font_name="Arial"):
    """Add a text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_multiline(slide, left, top, width, height, lines, font_size=10, color=PAPER_CANVAS, line_spacing=Pt(16)):
    """Add text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, bold, clr) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = clr
        p.font.name = "Arial"
        p.space_after = Pt(4)
    return txBox

def add_screenshot(slide, img_path, left, top, max_w, max_h):
    """Add screenshot image, cropped to top portion if tall, fit within max dimensions."""
    from PIL import Image as PILImage
    img = PILImage.open(img_path)
    iw, ih = img.size

    # For tall screenshots, crop to top portion (roughly 16:10 aspect)
    target_aspect = max_w / max_h
    current_aspect = iw / ih

    if current_aspect < target_aspect * 0.8:
        # Very tall image - crop to top portion
        crop_h = int(iw / target_aspect)
        if crop_h < ih:
            cropped = img.crop((0, 0, iw, crop_h))
            crop_path = img_path.replace(".png", "-crop.png")
            cropped.save(crop_path, quality=95)
            img_path = crop_path
            iw, ih = cropped.size

    # Scale to fit
    scale = min(max_w / iw, max_h / ih)
    final_w = int(iw * scale)
    final_h = int(ih * scale)

    pic = slide.shapes.add_picture(img_path, left, top, final_w, final_h)

    # Add subtle border
    pic.line.color.rgb = RGBColor(0x22, 0x58, 0x2C)
    pic.line.width = Pt(1.5)

    return pic, final_w, final_h

def add_badge(slide, left, top, width, height, text, bg_color, text_color=PAPER_CANVAS, font_size=8):
    """Add a rounded badge/tag."""
    shape = slide.shapes.add_shape(
        5,  # MSO_SHAPE.ROUNDED_RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = text_color
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return shape

def move_slide(prs, old_idx, new_idx):
    """Move a slide from old_idx to new_idx (0-based)."""
    slides = prs.slides._sldIdLst
    el = list(slides)[old_idx]
    slides.remove(el)
    if new_idx >= len(list(slides)):
        slides.append(el)
    else:
        ref = list(slides)[new_idx]
        slides.insert(slides.index(ref), el)

# ============================================================
# MAIN
# ============================================================
prs = Presentation(SRC)
blank_layout = prs.slide_layouts[0]  # Only layout available (DEFAULT)

# Get slide count before adding
original_count = len(prs.slides)
print(f"Original slides: {original_count}")

# ============================================================
# SLIDE A: After slide 4 (Design di Guglielmo) - "Il Design dal Vivo"
# Shows: Homepage + Servizio Progettazione
# ============================================================
slideA = prs.slides.add_slide(blank_layout)
add_bg(slideA, DEEP_FOREST)

# Title
add_text(slideA, Inches(0.55), Inches(0.3), Inches(9), Inches(0.5),
         "Il Design dal Vivo", font_size=22, bold=True, color=SUN_ACCENT)

# Subtitle
add_text(slideA, Inches(0.55), Inches(0.75), Inches(9), Inches(0.3),
         "La palette e i componenti Stitch applicati alle pagine reali del sito",
         font_size=9, bold=True, color=LIGHT_GREEN)

# Screenshot 1: Homepage
add_text(slideA, Inches(0.55), Inches(1.15), Inches(4), Inches(0.25),
         "Homepage", font_size=10, bold=True, color=PAPER_CANVAS)

pic1, w1, h1 = add_screenshot(slideA,
    f"{SCREENSHOTS}/00-homepage/desktop.png",
    Inches(0.55), Inches(1.4), Inches(4.0), Inches(2.9))

# Screenshot 2: Servizio Progettazione
add_text(slideA, Inches(5.1), Inches(1.15), Inches(4.5), Inches(0.25),
         "Servizio Progettazione Giardini", font_size=10, bold=True, color=PAPER_CANVAS)

pic2, w2, h2 = add_screenshot(slideA,
    f"{SCREENSHOTS}/10-servizio-progettazione-giardini/desktop.png",
    Inches(5.1), Inches(1.4), Inches(4.0), Inches(2.9))

# Annotation box
add_multiline(slideA, Inches(0.55), Inches(4.5), Inches(8.9), Inches(0.9), [
    ("Guglielmo (Design)", True, SUN_ACCENT),
    ("Palette High-Contrast Nature, Glass Navbar, card system e spacing definiti in Stitch", False, PAPER_CANVAS),
], font_size=9)

add_multiline(slideA, Inches(4.8), Inches(4.5), Inches(4.7), Inches(0.9), [
    ("Luca (Tech)", True, SUN_ACCENT),
    ("ServiceHero, ProcessSteps, VideoShowcase, AccordionFAQ implementati come componenti React", False, PAPER_CANVAS),
], font_size=9)

# ============================================================
# SLIDE B: After slide 5 (Processo Narrativo) - "La Voce di Barbara"
# Shows: Chi Siamo + Qualita
# ============================================================
slideB = prs.slides.add_slide(blank_layout)
add_bg(slideB, DEEP_FOREST)

add_text(slideB, Inches(0.55), Inches(0.3), Inches(9), Inches(0.5),
         "La Voce di Barbara sul Sito", font_size=22, bold=True, color=SUN_ACCENT)

add_text(slideB, Inches(0.55), Inches(0.75), Inches(9), Inches(0.3),
         "Il copy narrativo \"Il Verde che Vive\" applicato alle pagine Chi Siamo e Qualita'",
         font_size=9, bold=True, color=LIGHT_GREEN)

# Screenshot 1: Chi Siamo
add_text(slideB, Inches(0.55), Inches(1.15), Inches(4), Inches(0.25),
         "Chi Siamo", font_size=10, bold=True, color=PAPER_CANVAS)

pic3, w3, h3 = add_screenshot(slideB,
    f"{SCREENSHOTS}/01-chi-siamo/desktop.png",
    Inches(0.55), Inches(1.4), Inches(4.0), Inches(2.9))

# Screenshot 2: Qualita
add_text(slideB, Inches(5.1), Inches(1.15), Inches(4.5), Inches(0.25),
         "Qualita'", font_size=10, bold=True, color=PAPER_CANVAS)

pic4, w4, h4 = add_screenshot(slideB,
    f"{SCREENSHOTS}/06-qualita/desktop.png",
    Inches(5.1), Inches(1.4), Inches(4.0), Inches(2.9))

# Annotations
add_multiline(slideB, Inches(0.55), Inches(4.5), Inches(4.2), Inches(0.9), [
    ("Barbara (Storytelling)", True, SUN_ACCENT),
    ("Narrativa \"Il Verde che Vive\", seed content e rewrite completo per 4 pagine chiave", False, PAPER_CANVAS),
], font_size=9)

add_multiline(slideB, Inches(4.8), Inches(4.5), Inches(4.7), Inches(0.9), [
    ("Luca (Tech)", True, SUN_ACCENT),
    ("Copy integrato nei componenti React, filosofia e valori collegati ai servizi", False, PAPER_CANVAS),
], font_size=9)

# ============================================================
# SLIDE C: After slide 6 (Lead Generation) - "Il Quiz in Azione"
# Shows: Quiz Landing + Domanda + Risultato
# ============================================================
slideC = prs.slides.add_slide(blank_layout)
add_bg(slideC, DEEP_FOREST)

add_text(slideC, Inches(0.55), Inches(0.3), Inches(9), Inches(0.5),
         "Il Quiz in Azione", font_size=22, bold=True, color=SUN_ACCENT)

add_text(slideC, Inches(0.55), Inches(0.75), Inches(9), Inches(0.3),
         "6 domande, profilo giardiniere, scorecard personalizzata, lead capture",
         font_size=9, bold=True, color=LIGHT_GREEN)

# Screenshot 1: Quiz Landing
add_text(slideC, Inches(0.15), Inches(1.15), Inches(3), Inches(0.25),
         "Landing Quiz", font_size=9, bold=True, color=PAPER_CANVAS)

pic5, w5, h5 = add_screenshot(slideC,
    f"{SCREENSHOTS}/07-quiz-landing/desktop.png",
    Inches(0.15), Inches(1.4), Inches(3.05), Inches(2.9))

# Screenshot 2: Quiz Domanda
add_text(slideC, Inches(3.4), Inches(1.15), Inches(3), Inches(0.25),
         "Domanda 1", font_size=9, bold=True, color=PAPER_CANVAS)

pic6, w6, h6 = add_screenshot(slideC,
    f"{SCREENSHOTS}/50-quiz-domanda-1/desktop.png",
    Inches(3.4), Inches(1.4), Inches(3.05), Inches(2.9))

# Screenshot 3: Quiz Risultato
add_text(slideC, Inches(6.65), Inches(1.15), Inches(3), Inches(0.25),
         "Risultato Scorecard", font_size=9, bold=True, color=PAPER_CANVAS)

pic7, w7, h7 = add_screenshot(slideC,
    f"{SCREENSHOTS}/56-quiz-risultato-contemplativo/desktop.png",
    Inches(6.65), Inches(1.4), Inches(3.05), Inches(2.9))

# Annotations
add_multiline(slideC, Inches(0.15), Inches(4.5), Inches(4.8), Inches(0.9), [
    ("Luca (Tech)", True, SUN_ACCENT),
    ("Quiz micro-funnel con 6 domande, scoring engine, 4 profili giardiniere, schema Convex, QuizCTA widget", False, PAPER_CANVAS),
], font_size=9)

add_multiline(slideC, Inches(5.2), Inches(4.5), Inches(4.5), Inches(0.9), [
    ("Guglielmo (Design)", True, SUN_ACCENT),
    ("Template Stitch per ogni step del quiz, palette coerente, card risposte e scorecard layout", False, PAPER_CANVAS),
], font_size=9)

# ============================================================
# SLIDE D: After slide 8 (Stack) - "Le Pagine del Sito"
# Shows: Servizi + Progetti + Blog
# ============================================================
slideD = prs.slides.add_slide(blank_layout)
add_bg(slideD, DEEP_FOREST)

add_text(slideD, Inches(0.55), Inches(0.3), Inches(9), Inches(0.5),
         "Le Pagine del Sito", font_size=22, bold=True, color=SUN_ACCENT)

add_text(slideD, Inches(0.55), Inches(0.75), Inches(9), Inches(0.3),
         "Servizi categorizzati, portfolio progetti migrato a Convex, blog con articoli SEO",
         font_size=9, bold=True, color=LIGHT_GREEN)

# Screenshot 1: Servizi
add_text(slideD, Inches(0.15), Inches(1.15), Inches(3), Inches(0.25),
         "Servizi (12 pagine)", font_size=9, bold=True, color=PAPER_CANVAS)

pic8, w8, h8 = add_screenshot(slideD,
    f"{SCREENSHOTS}/02-servizi/desktop.png",
    Inches(0.15), Inches(1.4), Inches(3.05), Inches(2.9))

# Screenshot 2: Progetti
add_text(slideD, Inches(3.4), Inches(1.15), Inches(3), Inches(0.25),
         "Progetti (25+ case study)", font_size=9, bold=True, color=PAPER_CANVAS)

pic9, w9, h9 = add_screenshot(slideD,
    f"{SCREENSHOTS}/03-progetti/desktop.png",
    Inches(3.4), Inches(1.4), Inches(3.05), Inches(2.9))

# Screenshot 3: Blog
add_text(slideD, Inches(6.65), Inches(1.15), Inches(3), Inches(0.25),
         "Blog (3 articoli SEO)", font_size=9, bold=True, color=PAPER_CANVAS)

pic10, w10, h10 = add_screenshot(slideD,
    f"{SCREENSHOTS}/04-blog/desktop.png",
    Inches(6.65), Inches(1.4), Inches(3.05), Inches(2.9))

# Annotations
add_multiline(slideD, Inches(0.15), Inches(4.5), Inches(3.1), Inches(0.9), [
    ("Luca (Tech)", True, SUN_ACCENT),
    ("Migrazione blog e progetti da file statici a Convex DB, admin CRUD, dynamic routes", False, PAPER_CANVAS),
], font_size=9)

add_multiline(slideD, Inches(3.4), Inches(4.5), Inches(3.1), Inches(0.9), [
    ("Barbara (Storytelling)", True, SUN_ACCENT),
    ("Copy servizi, descrizioni progetti, articoli blog ottimizzati SEO", False, PAPER_CANVAS),
], font_size=9)

add_multiline(slideD, Inches(6.65), Inches(4.5), Inches(3.1), Inches(0.9), [
    ("Guglielmo (Design)", True, SUN_ACCENT),
    ("Card layout, filtri categorie, image grid, hero sections da template Stitch", False, PAPER_CANVAS),
], font_size=9)

# ============================================================
# REORDER: Move new slides to correct positions
# After appending, indices are:
#   0-9: original slides, 10:A, 11:B, 12:C, 13:D
# Target:
#   0:Title, 1:Dove, 2:Tre, 3:Design, 4:A, 5:Processo,
#   6:B, 7:Lead, 8:C, 9:Pronti, 10:Stack, 11:D, 12:Wave, 13:Closing
# ============================================================

# Step 1: Move A (idx 10) -> position 4
move_slide(prs, 10, 4)
# Now: 0,1,2,3,A,4,5,6,7,8,9,B,C,D

# Step 2: Move B (now at idx 11) -> position 6
move_slide(prs, 11, 6)
# Now: 0,1,2,3,A,4,B,5,6,7,8,9,C,D

# Step 3: Move C (now at idx 12) -> position 8
move_slide(prs, 12, 8)
# Now: 0,1,2,3,A,4,B,5,C,6,7,8,9,D

# Step 4: Move D (now at idx 13) -> position 11
move_slide(prs, 13, 11)
# Now: 0,1,2,3,A,4,B,5,C,6,7,D,8,9

prs.save(OUT)
print(f"Saved enriched presentation to {OUT}")
print(f"Total slides: {len(prs.slides)}")
